"""Build the 24-minute talk deck.

Two rules govern the split between what is on a slide and what is in the
notes, and they are the whole design:

**A slide is looked at, not read.** Anything longer than a phrase competes
with the speaker — the room reads it, stops listening, and finishes before
the sentence does. So slides carry a number, a name, or four short labels.
Never an argument.

**The notes carry the argument.** They are written as a presenter's
guideline, not a transcript: what to say, the one thing to emphasise, what
to point at, and the question that is coming. Structured so a glance at the
confidence monitor finds the next beat rather than a paragraph to re-read.

Dark, because every screenshot and clip in it is dark: a light deck wrapped
around dark product shots reads as two decks stapled together.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SHOTS = ROOT / "docs" / "slide-screenshots"
DEMOS = ROOT / "docs" / "demos"
OUT = ROOT / "docs" / "reading-the-paperwork.pptx"

# Product dark palette — the same values the app ships.
BG      = RGBColor(0x0B, 0x0F, 0x0E)
SURFACE = RGBColor(0x16, 0x1C, 0x1B)
TEXT    = RGBColor(0xF2, 0xF5, 0xF4)
MUTED   = RGBColor(0x98, 0xA3, 0xA1)
DIM     = RGBColor(0x6D, 0x78, 0x77)
ACCENT  = RGBColor(0xA7, 0x9B, 0xFF)
AMBER   = RGBColor(0xDE, 0x75, 0x29)
RULE    = RGBColor(0x2A, 0x33, 0x32)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════
# presenter notes
# ══════════════════════════════════════════════════════════════════════

def guide(*, at=None, dur=None, say=(), emphasise=None, point=(), asked=(),
          nxt=None, warn=None):
    """Format one slide's notes as a guideline rather than a script.

    Fixed section order, so the eye lands in the same place every time: the
    beats, then the single thing to stress, then what to point at, then the
    question that is coming. A presenter glancing down mid-sentence needs
    to find their place in about a second.
    """
    out = []
    if at:
        out += [f"══  {at}   ·   {dur}  ══", ""]
    if say:
        out += ["SAY"] + [f"  •  {s}" for s in say] + [""]
    if emphasise:
        out += ["THE ONE THING", f"  ►  {emphasise}", ""]
    if point:
        out += ["POINT AT"] + [f"  ·  {p}" for p in point] + [""]
    if warn:
        out += ["WATCH OUT", f"  !  {warn}", ""]
    if asked:
        out += ["IF ASKED"]
        for q, a in asked:
            out += [f"  Q  {q}", f"  A  {a}"]
        out += [""]
    if nxt:
        out += [f"NEXT  →  {nxt}"]
    return "\n".join(out).rstrip()


def new_slide(notes: str = ""):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def text(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (string, size_pt, bold, colour, space_after_pt)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (s, size, bold, colour, after) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.alignment = align
        p.space_after = Pt(after)
        f = p.font
        f.size, f.bold, f.color.rgb = Pt(size), bold, colour
        f.name = "Segoe UI"
    return tb


def eyebrow_title(slide, eyebrow, title):
    text(slide, MARGIN, Inches(0.62), CONTENT_W, Inches(0.3),
         [(eyebrow.upper(), 12, True, ACCENT, 0)])
    text(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(1.2),
         [(title, 36, True, TEXT, 0)])


def timing(slide, clock, dur, idx):
    text(slide, W - Inches(2.3), Inches(0.62), Inches(1.5), Inches(0.6),
         [(f"{clock}  ·  {dur}", 11, True, DIM, 0),
          (f"slide {idx}", 10, False, DIM, 0)], align=PP_ALIGN.RIGHT)


def labels(slide, items, top=Inches(2.9), size=24, gap=20):
    """Short labels, not sentences. If one wraps, it is too long."""
    runs = [(it, size, False, TEXT, gap) for it in items]
    text(slide, MARGIN, top, CONTENT_W * 0.9, Inches(4.0), runs)


def hero(slide, big, under=None, top=Inches(2.6), size=96):
    text(slide, MARGIN, top, CONTENT_W, Inches(1.9), [(big, size, True, ACCENT, 6)])
    if under:
        text(slide, MARGIN, top + Inches(1.75), CONTENT_W * 0.8, Inches(1.0),
             [(under, 26, False, TEXT, 0)])


def shot(slide, name, top=Inches(2.35), height=Inches(4.5)):
    path = SHOTS / f"{name}.png"
    if not path.exists():
        return
    from PIL import Image
    with Image.open(path) as im:
        ratio = im.width / im.height
    width = Emu(int(height * ratio))
    left = int((W - width) / 2)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left - Inches(0.06), top - Inches(0.06),
                                   width + Inches(0.12), height + Inches(0.12))
    frame.fill.solid()
    frame.fill.fore_color.rgb = SURFACE
    frame.line.color.rgb = RULE
    frame.line.width = Pt(1)
    frame.shadow.inherit = False
    slide.shapes.add_picture(str(path), left, top, height=height)


def video(slide, name, left, top, width):
    """Embed a clip, with its poster frame ALSO laid in as a real picture.

    Canva's .pptx import does not carry embedded media across. Relying on
    the poster inside the media shape would mean relying on that shape
    surviving, which is the thing that does not. A separate picture costs a
    few MB and means the deck degrades to a screenshot deck rather than to
    blank boxes — and the clip drops back on top in one action.
    """
    clip = DEMOS / f"{name}.mp4"
    poster = DEMOS / "posters" / f"{name}.png"
    if not clip.exists():
        return
    height = Emu(int(width / 1.6))  # recordings are 1440x900
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left - Inches(0.06), top - Inches(0.06),
                                   width + Inches(0.12), height + Inches(0.12))
    frame.fill.solid()
    frame.fill.fore_color.rgb = SURFACE
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(1.25)
    frame.shadow.inherit = False
    if poster.exists():
        slide.shapes.add_picture(str(poster), left, top, width, height)
    slide.shapes.add_movie(str(clip), left, top, width, height,
                           poster_frame_image=str(poster) if poster.exists() else None,
                           mime_type="video/mp4")


def clip_slide(name, eyebrow, title, watch, notes, clock=None, dur=None, idx=None):
    """Clip on the right; on the left only what to watch for, in phrases."""
    s = new_slide(notes)
    eyebrow_title(s, eyebrow, title)
    if clock:
        timing(s, clock, dur, idx)
    text(s, MARGIN, Inches(2.55), Inches(4.0), Inches(4.2),
         [(w, 17, False, TEXT, 16) for w in watch])
    video(s, name, Inches(5.25), Inches(2.45), Inches(7.15))
    return s


def card(slide, left, top, width, height, head, body=None, colour=None):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = colour or RULE
    b.line.width = Pt(1.4)
    b.shadow.inherit = False
    runs = [(head, 20, True, colour or TEXT, 8)]
    if body:
        runs.append((body, 16, False, MUTED, 0))
    text(slide, left + Inches(0.4), top + Inches(0.3),
         width - Inches(0.8), height - Inches(0.6), runs)


# ══════════════════════════════════════════════════════════════════════
# 01 — Title
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="0:00", dur="0:20",
    say=["Name, and where the internship was.",
         "One sentence on what the next 24 minutes are for."],
    emphasise="Do not narrate the agenda. An agenda slide spends 30 seconds "
              "telling people they will later learn something.",
    nxt="02 — the cost"))
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.15), Inches(1.5), Pt(5))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
bar.shadow.inherit = False
text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.4),
     [("END-OF-STUDIES INTERNSHIP  ·  SOPRA HR SOFTWARE × ESPRIT", 13, True, ACCENT, 0)])
text(s, MARGIN, Inches(2.55), CONTENT_W, Inches(1.6),
     [("Reading the Paperwork", 60, True, TEXT, 0)])
text(s, MARGIN, Inches(4.05), Inches(8.6), Inches(1.0),
     [("HR documents → machine-learning features", 22, False, MUTED, 0)])
text(s, MARGIN, Inches(6.2), CONTENT_W, Inches(0.9),
     [("[Your name]", 17, True, TEXT, 4),
      ("Supervisors: [company] · [academic]      [date]", 14, False, DIM, 0)])

# ══════════════════════════════════════════════════════════════════════
# 02 — The cost
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="0:20", dur="1:00",
    say=["One resignation costs recruitment, onboarding, lost productivity.",
         "And the knowledge that left, which nobody wrote down.",
         "Expensive, consequential — and partly predictable.",
         "That is why every HR platform wants a turnover model."],
    emphasise="Open on cost, not technology. The commercial half of the room "
              "needs a reason to listen before any architecture appears.",
    warn="FILL THIS IN with a sourced figure. An invented number is the first "
         "thing a jury challenges — and the only number on the slide.",
    asked=[("Where does that figure come from?",
            "Name the source out loud. If you are unsure of it, say the range "
            "rather than the point estimate.")],
    nxt="03 — but the data is unreadable"))
eyebrow_title(s, "Act I · the problem worth money", "What one resignation costs")
timing(s, "0:20", "1:00", "02")
hero(s, "[ X ]×", "monthly salary, to replace one person")
text(s, MARGIN, Inches(6.3), CONTENT_W, Inches(0.5),
     [("↑  replace with a sourced figure", 14, False, AMBER, 0)])

# ══════════════════════════════════════════════════════════════════════
# 03 — The unread archive
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="1:20", dur="1:20",
    say=["Here is what surprised me. The obstacle is not modelling.",
         "Published work runs on a few public datasets. The most-used one has "
         "about 1,470 rows — and is itself synthetic.",
         "Meanwhile a twenty-year-old company holds thousands of reviews, every "
         "resignation letter it ever received, and a great deal of paper.",
         "None of it reaches any model."],
    emphasise="Say the last line slowly. It is the thesis of the whole talk, "
              "and everything after it is a consequence.",
    point=["The three archive lines, then the quote"],
    nxt="04 — what already existed"))
eyebrow_title(s, "Act I", "The data exists")
timing(s, "1:20", "1:20", "03")
labels(s, ["Thousands of performance reviews",
           "Every resignation letter ever received",
           "Offer letters, CVs — and a lot of paper"], top=Inches(2.5), size=25)
text(s, MARGIN, Inches(5.5), CONTENT_W * 0.85, Inches(1.2),
     [("It just isn't readable.", 40, True, ACCENT, 0)])

# ══════════════════════════════════════════════════════════════════════
# 04 — What existed
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="2:40", dur="1:20",
    say=["The platform simulates an organisation week by week, so employees "
         "actually resign inside the model.",
         "That means the classifier trains on real events rather than a risk "
         "score somebody assigned.",
         "It worked. It had one problem: everything in it came from its own "
         "simulation. Nothing real ever entered."],
    emphasise="The difference between learning something and restating an "
              "assumption.",
    warn="80 seconds. Resist explaining the simulation — it is context, not "
         "your contribution. If you are running long, this is the first cut.",
    nxt="05 — the gap"))
eyebrow_title(s, "Act I", "What already existed")
timing(s, "2:40", "1:20", "04")
shot(s, "atrisk", top=Inches(2.2), height=Inches(4.5))

# ══════════════════════════════════════════════════════════════════════
# 05 — The gap
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="4:00", dur="1:20",
    say=["When the scoring layer needed performance history, it had none.",
         "So it supplied a neutral 3 out of 5 — for everybody.",
         "Three of eighteen features were the same number for every employee. "
         "A tree never splits on them, so the effective count was fifteen.",
         "And a declining rating is one of the more reliable early signals of "
         "a resignation. The missing signal was not an arbitrary one."],
    emphasise="This is the pivot. From here the audience knows exactly what "
              "problem the next fifteen minutes solve.",
    point=["The three constants, then the 3-of-18"],
    nxt="06 — the pipeline"))
eyebrow_title(s, "Act I", "Three features were constants")
timing(s, "4:00", "1:20", "05")
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(2.7), Inches(6.0), Inches(2.1))
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
box.line.color.rgb = RULE
box.shadow.inherit = False
text(s, MARGIN + Inches(0.45), Inches(3.0), Inches(5.2), Inches(1.6),
     [("rating_last   =  3.0", 23, True, TEXT, 10),
      ("rating_prev   =  3.0", 23, True, TEXT, 10),
      ("rating_delta  =  0.0", 23, True, TEXT, 0)])
text(s, Inches(7.5), Inches(2.85), Inches(5.0), Inches(2.2),
     [("3 of 18", 56, True, AMBER, 8),
      ("carried no information", 20, False, TEXT, 0)])

# ══════════════════════════════════════════════════════════════════════
# 06 — The pipeline
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="5:20", dur="1:40",
    say=["Documents arrive as CSV, PDF, free text — or a photograph.",
         "Structured files parse deterministically. A roster needs no model and "
         "no API key at all.",
         "Prose goes to a language model. Paper goes through OCR.",
         "All three routes end in the same place: a staging table.",
         "There is no code path from extraction to a personnel record."],
    emphasise="'No unreviewed write' is a property of the architecture, not a "
              "rule someone follows.",
    point=["The STAGING box — it is the largest for a reason",
           "Then the HUMAN box"],
    warn="You will point back at this slide three times. Make sure the room "
         "sees the staging box now.",
    nxt="07 — the demo"))
eyebrow_title(s, "Act II · what I built", "One pipeline, one gate")
timing(s, "5:20", "1:40", "06")
stages = [("Document\nCSV · PDF · text · PHOTO", SURFACE, TEXT, 2.10),
          ("Parse\nrules · LLM · OCR", SURFACE, TEXT, 2.10),
          ("STAGING\nnothing proceeds", ACCENT, BG, 2.45),
          ("HUMAN\napproves per field", SURFACE, TEXT, 2.10),
          ("Employees · reviews\nexit notes · cohort", SURFACE, TEXT, 2.15)]
x = MARGIN
for label, fill, fg, w_in in stages:
    w = Inches(w_in)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.1), w, Inches(1.6))
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = ACCENT if fill == ACCENT else RULE
    b.line.width = Pt(1.5)
    b.shadow.inherit = False
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15 if i == 0 else 11.5)
        p.font.bold = i == 0
        p.font.color.rgb = fg
        p.font.name = "Segoe UI"
    x += w + Inches(0.12)
text(s, MARGIN, Inches(5.3), CONTENT_W, Inches(0.7),
     [("No path from extraction to a personnel record.", 22, True, ACCENT, 0)])

# ══════════════════════════════════════════════════════════════════════
# 07 — Demo
# ══════════════════════════════════════════════════════════════════════
clip_slide(
    "03-photographed-paper-ocr", "Act II", "A photograph of paper",
    ["Uploaded → marked OCR",
     "Transcribed once, on upload",
     "Then the identical route",
     "Staged at 0.7× confidence"],
    clock="7:00", dur="2:30", idx="07",
    notes=guide(
        at="7:00", dur="2:30",
        say=["Narrate what is happening, not what you are clicking.",
             "A photograph produces the same text a typed document would.",
             "From there, nothing downstream is special-cased for paper.",
             "Only the confidence differs — and it differs automatically."],
        emphasise="Paper is a front door, not a second pipeline.",
        point=["The OCR badge on the row",
               "The confidence figure when the facts stage"],
        warn="This is a recording, on purpose. A live cloud call on conference "
             "wifi is a two-minute silence. Let it play; do not talk over the "
             "whole clip — narrate the first beat, then stop.",
        asked=[("How accurate is the transcription?",
                "Honest answer: not measured against a real archive. That is "
                "why every fact from a photo is discounted rather than trusted."),
               ("Why not just use Tesseract?",
                "It is one of the four backends. A vision model reads for "
                "meaning, so it copes with skew and handwriting.")],
        nxt="08 — principle 1"))

# ══════════════════════════════════════════════════════════════════════
# 08 — Principle 1
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="9:30", dur="1:20",
    say=["Both parsers write only to the staging table.",
         "The apply endpoint is the only writer of employee data, and it acts "
         "only on facts approved by identifier.",
         "Every staged change carries its current value, the proposed value, a "
         "confidence — and the sentence it came from."],
    emphasise="For a vendor holding client personnel data under GDPR this is "
              "not a nice-to-have. It is the difference between a feature you "
              "can sell and one legal will not let you ship.",
    point=["The evidence column — approving is a decision, not an act of faith"],
    warn="Say the GDPR line looking at the business half of the room. This is "
         "where a sceptical buyer relaxes.",
    nxt="09 — principle 2"))
eyebrow_title(s, "Act III · why it can be trusted", "Nothing is written without a human")
timing(s, "9:30", "1:20", "08")
shot(s, "documents", top=Inches(2.2), height=Inches(4.5))

# ══════════════════════════════════════════════════════════════════════
# 09 — Principle 2
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="10:50", dur="2:00",
    say=["A resignation letter says the author was exhausted and unsupported.",
         "Workload and manager support are legitimate features. So why not use "
         "them?",
         "Because the letter is written AT the moment of the outcome.",
         "A model trained on it learns that people who write about exhaustion "
         "in resignation letters resign. Excellent AUC. Zero predictive value — "
         "at scoring time no such letter exists.",
         "So the schema has nowhere to put those values."],
    emphasise="A confused model, a careless contributor, a bad prompt — none "
              "of them can leak through this path. The field is not there.",
    point=["The four fields", "Then the two lines under them"],
    warn="Two full minutes. This is the strongest technical idea in the talk. "
         "If a jury remembers one thing, make it this.",
    asked=[("Couldn't you just filter those fields out later?",
            "You could, and it would work until someone forgot. A check that "
            "must be remembered eventually is not."),
           ("Isn't the note text itself a feature?",
            "It feeds sentiment analysis for humans to read, not the model's "
            "feature frame.")],
    nxt="10 — principle 3"))
eyebrow_title(s, "Act III", "Make the wrong thing unrepresentable")
timing(s, "10:50", "2:00", "09")
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(2.6), Inches(6.2), Inches(2.9))
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
box.line.color.rgb = RULE
box.shadow.inherit = False
text(s, MARGIN + Inches(0.45), Inches(2.9), Inches(5.4), Inches(2.4),
     [("ResignationLetterExtract", 18, True, ACCENT, 12),
      ("employee_email", 17, False, TEXT, 7),
      ("effective_date", 17, False, TEXT, 7),
      ("is_voluntary", 17, False, TEXT, 7),
      ("note_text", 17, False, TEXT, 0)])
text(s, Inches(7.7), Inches(2.9), Inches(4.8), Inches(3.0),
     [("No workload field.", 27, True, TEXT, 10),
      ("No morale field.", 27, True, TEXT, 18),
      ("No feature fields at all.", 24, True, AMBER, 0)])

# ══════════════════════════════════════════════════════════════════════
# 10 — Principle 3
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="12:50", dur="1:20",
    say=["An organisation never receives a resignation letter from someone who "
         "stayed.",
         "So a training set built from letters alone is a hundred per cent "
         "positive. The model learns that everybody leaves.",
         "Cohort construction refuses to build a training set unless a roster "
         "establishes who did not leave — and refuses outright above a fifty "
         "per cent base rate."],
    emphasise="The roster is not an optional enrichment. It is what makes the "
              "labels mean anything at all.",
    point=["100% on the left, then 5.99% on the right"],
    warn="The most accessible idea in the talk — non-technical listeners "
         "follow it completely. If you must cut slide 09, this one survives.",
    nxt="11 — when the model lied"))
eyebrow_title(s, "Act III", "You only get letters from leavers")
timing(s, "12:50", "1:20", "10")
for title, sub, val, colour, lx in [
        ("Letters only", "24 positives · 0 negatives", "100%", AMBER, 0.9),
        ("Letters + roster", "24 positives · 377 negatives", "5.99%", ACCENT, 7.1)]:
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(2.8),
                           Inches(5.35), Inches(2.9))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = colour
    b.line.width = Pt(1.6)
    b.shadow.inherit = False
    text(s, Inches(lx + 0.45), Inches(3.05), Inches(4.5), Inches(2.4),
         [(title, 19, True, TEXT, 6), (sub, 15, False, MUTED, 16),
          (val, 60, True, colour, 4), ("base rate", 14, False, DIM, 0)])

# ══════════════════════════════════════════════════════════════════════
# 11 — The model made something up
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="14:10", dur="1:20",
    say=["My favourite failure.",
         "The prompt told the model to report a missing field. It complied "
         "semantically and failed structurally — it put the refusal INSIDE the "
         "field.",
         "Schema validation passed, because that is a perfectly valid string.",
         "It staged a plausible-looking hire.",
         "The fix was not a better prompt. The department is now resolved "
         "against departments that actually exist."],
    emphasise="Prompt wording constrains one phrasing of a failure. A "
              "structural check constrains the whole class of it.",
    point=["The string, then the line under it"],
    warn="Tell this as an anecdote — it is the one place in the talk to sound "
         "amused. It does more work than any results slide, because it shows "
         "you tested adversarially.",
    asked=[("Would a better model have avoided it?",
            "Maybe, on that phrasing. That is exactly why the fix is not a "
            "prompt.")],
    nxt="12 — paper and portability"))
eyebrow_title(s, "Act III", "When the model made something up")
timing(s, "14:10", "1:20", "11")
text(s, MARGIN, Inches(2.5), CONTENT_W * 0.9, Inches(0.6),
     [("The prompt said: if the department is missing, return an error.", 19, False, MUTED, 0)])
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(3.25), Inches(10.4), Inches(1.35))
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
box.line.color.rgb = AMBER
box.line.width = Pt(1.6)
box.shadow.inherit = False
text(s, MARGIN + Inches(0.5), Inches(3.5), Inches(9.4), Inches(0.9),
     [('department_name:  "no department stated"', 26, True, AMBER, 0)])
text(s, MARGIN, Inches(5.1), CONTENT_W * 0.9, Inches(1.4),
     [("A valid string. It passed validation.", 26, True, TEXT, 12),
      ("Structure catches a class. Wording catches a phrasing.", 20, False, ACCENT, 0)])

# ══════════════════════════════════════════════════════════════════════
# 12 — Paper and portability
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="15:30", dur="1:20",
    say=["Plenty of HR records are still paper.",
         "A photograph produces the same text a typed document would, and "
         "travels the identical route.",
         "What differs is confidence — and it differs automatically. OCR can "
         "turn a three into an eight in a salary, and nothing downstream would "
         "catch it, because both are valid numbers.",
         "And the provider is a configuration setting, not a code decision."],
    emphasise="A client who requires everything inside their own AWS account "
              "gets that without a code change.",
    point=["0.7× on the left", "One environment variable on the right"],
    warn="Two commercial objections answered before they are asked: 'our "
         "clients aren't digitised' and 'we can't send data to a third party'.",
    asked=[("What does OCR cost per page?",
            "Metered separately under its own feature — roughly a thousand "
            "tokens for a typical page.")],
    nxt="13 — results"))
eyebrow_title(s, "Act III", "Paper, and not being locked in")
timing(s, "15:30", "1:20", "12")
card(s, Inches(0.9), Inches(2.7), Inches(5.35), Inches(3.0), "Paper",
     "Photo → transcription → the same pipeline.\n\nStaged at 0.7× confidence.\n\n"
     "A scanned PDF takes the same route.", ACCENT)
card(s, Inches(7.1), Inches(2.7), Inches(5.35), Inches(3.0), "Portability",
     "Groq ⇄ AWS Bedrock.\n\nOne environment variable.\n\n"
     "Runs inside the client's own account.", ACCENT)

# ══════════════════════════════════════════════════════════════════════
# 13 — Results
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="16:50", dur="1:30",
    say=["I generated documents across six types.",
         "About one in six was built to be REFUSED — a review with no rating, "
         "a redundancy notice, a letter from someone who does not work here.",
         "A corpus of happy paths proves nothing about whether a parser invents "
         "data when a document is deficient.",
         "All seventeen refusals behaved correctly, and with specific reasons "
         "rather than a generic failure."],
    emphasise="Say 'designed to fail' out loud. Testing the refusals is what "
              "distinguishes engineering from a demo.",
    point=["17 / 17 first, then the other three"],
    asked=[("Why generated documents rather than real ones?",
            "Real HR documents are personal data. It is a genuine limitation "
            "and it is on slide 16.")],
    nxt="14 — the second corpus"))
eyebrow_title(s, "Act IV · what it produced", "Tested against documents built to fail")
timing(s, "16:50", "1:30", "13")
for i, (val, lab) in enumerate([("17/17", "refusals correct"),
                                ("143", "documents"),
                                ("350", "tests"),
                                ("6", "formats")]):
    lx = Inches(0.9 + i * 3.05)
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, Inches(3.0), Inches(2.85), Inches(2.2))
    b.fill.solid()
    b.fill.fore_color.rgb = SURFACE
    b.line.color.rgb = RULE
    b.shadow.inherit = False
    text(s, lx + Inches(0.3), Inches(3.35), Inches(2.3), Inches(1.7),
         [(val, 46, True, ACCENT, 8), (lab, 16, False, MUTED, 0)])

# ══════════════════════════════════════════════════════════════════════
# 14 — The second corpus
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="18:20", dur="1:20",
    say=["Those documents all tested CONTENT.",
         "So I built a second corpus that tests the CONTAINER — delimiters, "
         "encodings, locales, scales, page layout.",
         "It found this before a single file was uploaded.",
         "The parser strips every comma as a thousands separator. So a French "
         "salary is read a hundred times too large.",
         "Worse than a refusal — the wrong number reaches a reviewer looking "
         "completely plausible."],
    emphasise="Every one of my first 106 documents was UTF-8, comma-delimited "
              "and in English, because that is what you write when you are "
              "inventing plausible documents.",
    point=["The two rows, then the line under them"],
    warn="Say plainly that it is still open. For an HR vendor with "
         "French-locale clients this is commercial, not just technical.",
    nxt="15 — the gap closed"))
eyebrow_title(s, "Act IV", "The bug the first corpus could not find")
timing(s, "18:20", "1:20", "14")
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(2.6), Inches(10.4), Inches(1.9))
box.fill.solid()
box.fill.fore_color.rgb = SURFACE
box.line.color.rgb = AMBER
box.line.width = Pt(1.6)
box.shadow.inherit = False
text(s, MARGIN + Inches(0.55), Inches(2.9), Inches(9.2), Inches(1.4),
     [("1234,56        →        123456.0", 28, True, TEXT, 12),
      ("85.000,50      →        85.0005", 28, True, TEXT, 0)])
text(s, MARGIN, Inches(4.9), CONTENT_W * 0.92, Inches(1.6),
     [("A decimal comma multiplies salary by 100.", 28, True, AMBER, 10),
      ("Silently.", 28, True, AMBER, 0)])

# ══════════════════════════════════════════════════════════════════════
# 15 — The gap closed
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="19:40", dur="1:20",
    say=["Twelve employees now have two-period histories, with real deltas "
         "instead of a hard-coded zero.",
         "The promotion gate records document-sourced examples separately, so a "
         "change in AUC can be attributed to a source.",
         "And the fixed synthetic evaluation cohort never sees ingested data — "
         "so documents can improve the model or do nothing, but cannot quietly "
         "degrade production."],
    emphasise="Call back to slide 5 explicitly: 'those three constants'. "
              "Closing a loop the audience remembers opening is worth more "
              "than a new fact.",
    point=["The last row — 0 → 12"],
    nxt="16 — what doesn't work"))
eyebrow_title(s, "Act IV", "The gap from slide 5, closed")
timing(s, "19:40", "1:20", "15")
rows = [("Informative features", "15 of 18", "18 of 18"),
        ("Reviews in the database", "0", "38"),
        ("Employees with a history", "0", "26"),
        ("…with two periods", "0", "12")]
top = 2.9
for label, before, after in rows:
    text(s, MARGIN, Inches(top), Inches(6.4), Inches(0.5), [(label, 21, False, TEXT, 0)])
    text(s, Inches(7.9), Inches(top), Inches(1.6), Inches(0.5),
         [(before, 21, False, DIM, 0)], align=PP_ALIGN.RIGHT)
    text(s, Inches(9.8), Inches(top), Inches(0.6), Inches(0.5),
         [("→", 21, False, DIM, 0)], align=PP_ALIGN.CENTER)
    text(s, Inches(10.5), Inches(top), Inches(1.9), Inches(0.5),
         [(after, 21, True, ACCENT, 0)], align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(top + 0.55), CONTENT_W, Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    top += 0.85

# ══════════════════════════════════════════════════════════════════════
# 16 — What doesn't work
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="21:00", dur="1:00",
    say=["Three things I would not claim.",
         "The sentiment analysis is a small hand-built lexicon and it fails on "
         "real prose. A letter saying the author was tired in a way a holiday "
         "does not fix scores zero. At twenty-four documents, that is most of "
         "them.",
         "The decimal comma is still open — it has a fixture, not a fix.",
         "And I have not shown an improvement in predicting real turnover."],
    emphasise="That last one needs a real organisation's documents and a real "
              "outcome window. Neither was available.",
    warn="DO NOT SKIP THIS SLIDE to save time. Naming your own limits before a "
         "jury does is the cheapest credibility you will buy all talk — and it "
         "pre-empts the obvious hostile question.",
    asked=[("So does it actually work?",
            "It demonstrably ingests documents safely. Whether it improves "
            "prediction is unmeasured, and I am not claiming it.")],
    nxt="17 — the commercial case"))
eyebrow_title(s, "Act IV", "What doesn't work yet")
timing(s, "21:00", "1:00", "16")
labels(s, ["83% of exit notes yield no theme",
           "A decimal comma still corrupts salaries",
           "Duplicate emails can match the wrong person",
           "Validated on generated documents"], top=Inches(2.6), size=23, gap=18)
text(s, MARGIN, Inches(6.0), CONTENT_W * 0.9, Inches(0.8),
     [("No improvement in real turnover prediction is claimed.", 20, True, AMBER, 0)])

# ══════════════════════════════════════════════════════════════════════
# 17 — Commercial
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="22:00", dur="1:20",
    say=["The pitch in one line: this turns an archive a client is already "
         "paying to store into something their analytics can use.",
         "With an audit trail that survives a data-protection review.",
         "And because the cost per document is measured rather than estimated, "
         "a client can be quoted before a pilot rather than after it."],
    emphasise="Four claims, no more. Each answers a question a buyer actually "
              "asks: what do I get, can I defend it, where does my data go, "
              "what does it cost.",
    nxt="18 — close"))
eyebrow_title(s, "Act V · where it goes", "Why a client would pay for this")
timing(s, "22:00", "1:20", "17")
claims = [("Dormant asset", "Archives they already own"),
          ("Auditable", "Every value traces to a sentence"),
          ("Their infrastructure", "Runs in the client's own account"),
          ("Known unit cost", "~460 tokens, ~5 seconds per document")]
for i, (head, body) in enumerate(claims):
    lx = Inches(0.9 + (i % 2) * 6.2)
    ty = Inches(2.8 + (i // 2) * 1.9)
    card(s, lx, ty, Inches(5.5), Inches(1.55), head, body, ACCENT)

# ══════════════════════════════════════════════════════════════════════
# 18 — Close
# ══════════════════════════════════════════════════════════════════════
s = new_slide(guide(
    at="23:20", dur="0:40",
    say=["The highest-return next step is the sentiment model, because "
         "eighty-three per cent of what we ingest currently produces no signal.",
         "And if you take one idea away —",
         "Every safety property in this pipeline is structural. A letter cannot "
         "leak a feature because the schema has no feature fields. A batch of "
         "letters cannot become a training set because the base rate is checked.",
         "Each could have been a rule in a document. And each would eventually "
         "have been forgotten."],
    emphasise="Stop talking after that sentence. Leave the line on screen and "
              "let the silence be the invitation — do not add 'any questions?'",
    nxt="— questions —"))
eyebrow_title(s, "Act V", "Next, and close")
timing(s, "23:20", "0:40", "18")
text(s, MARGIN, Inches(2.5), CONTENT_W * 0.9, Inches(0.8),
     [("Next  —  the sentiment model  ·  the decimal comma  ·  a real archive",
       18, False, MUTED, 0)])
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(3.7), Inches(2.2), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
bar.shadow.inherit = False
text(s, MARGIN, Inches(4.2), CONTENT_W * 0.9, Inches(2.6),
     [("Make the wrong thing", 44, True, TEXT, 4),
      ("unrepresentable,", 44, True, TEXT, 4),
      ("rather than detectable.", 44, True, ACCENT, 0)])

# ══════════════════════════════════════════════════════════════════════
# Appendix — a clip per capability, outside the 24-minute run
# ══════════════════════════════════════════════════════════════════════
CLIPS = [
    ("01-roster-csv-deterministic-parsing", "Roster CSV",
     ["No model, no API key", "Compared row by row", "Only differences staged",
      "An identical row proposes nothing"]),
    ("02-resignation-letter-llm-extraction", "Resignation letter",
     ["Free prose, read by the model", "Contributes a label and a date",
      "Never a feature", "Note text flows to exit-note analysis"]),
    ("04-honest-refusal-of-an-involuntary-exit", "Honest refusal",
     ["A layoff is not a resignation", "Refused, with the reason stored",
      "Nothing was written", "That is the success case"]),
    ("05-review-queue-filter-and-evidence", "The review queue",
     ["142 documents", "Filter to what needs you", "List and inspector together",
      "Every proposal carries its sentence"]),
    ("06-scenario-simulator-forecast", "Scenario simulator",
     ["A template drops events", "Time is drawn, not tabulated",
      "Baseline against scenario", "Measured in the targeted cohort"]),
    ("07-diagnosis-report-modal", "Diagnosis report",
     ["Opens over the page", "Top drivers and cohort",
      "A recommendation you can apply", "Dismissing costs nothing"]),
    ("08-retention-risk-ranking", "Retention risk",
     ["Who ranks highest", "Each score traces to drivers", "Never a number alone"]),
    ("09-turnover-model-and-promotion-gate", "Model and promotion gate",
     ["Promoted only if AUC holds", "Every decision logged",
      "Document examples counted separately"]),
    ("10-token-usage-meter", "Token usage",
     ["Total, today, this week", "Split by feature",
      "The model that actually served it"]),
]

for _name, _title, _watch in CLIPS:
    clip_slide(_name, "Backup · not in the run", _title, _watch,
               notes=guide(say=["Backup clip. Not part of the 24-minute run.",
                                "Use only if the room asks for it."]))

for _name, _title in [("dashboard", "Dashboard"),
                      ("exitnotes", "Exit Notes Insights")]:
    _s = new_slide(guide(say=["Backup screen. Not part of the 24-minute run."]))
    eyebrow_title(_s, "Backup · not in the run", _title)
    shot(_s, _name, top=Inches(2.2), height=Inches(4.5))

prs.save(OUT)
print(f"wrote {OUT}")
print(f"  {len(prs.slides._sldIdLst)} slides")
